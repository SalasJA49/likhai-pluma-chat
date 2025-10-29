import os
from io import BytesIO
from urllib.parse import urlencode
from datetime import datetime

import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# PDF (install: pip install reportlab)
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# your modules
import app.pages as pages
import app.utils as utils
import app.prompts as prompts

from pages import reader as reader_page
from pages import outputs as outputs_page
from pages import settings as settings_page

st.set_page_config(page_title="Style Suite", layout="wide")
st.markdown("""
<style>
[data-testid="stSidebar"] .block-container{
  padding: -8px -12px 20px -12px !important;  /* top, right, bottom, left */
}

/* Center + size the logo */
[data-testid="stSidebar"] img.pluma-logo{
  display:block;
  margin: 0 auto;            /* center */
  width: 170px;              /* set your size here */
  max-width: 170px;
}

/* Pull the logo up a bit */
[data-testid="stSidebar"] .pluma-top{
  margin-top: -14px;         /* adjust: -8 to -24 looks good */
  margin-bottom: 16px;
}

/* Title & caption */
[data-testid="stSidebar"] .pluma-title{
  text-align:center; font-size:28px; font-weight:800; line-height:1.1; margin: 6px 0 4px;
}
[data-testid="stSidebar"] .pluma-caption{
  text-align:center; color:#6B7280; font-size:13.5px; line-height:1.35; margin-top:2px;
}

/* Optional: ensure logo is centered nicely */
[data-testid="stSidebar"] img{ display:block; margin-inline:auto; padding: -3px }

/* 1) Fix sidebar width & inner padding */
[data-testid="stSidebar"]{
  width: 300px !important; min-width:300px !important;
}

/* 2) Card-like expanders (Chatbot / Writer) */
[data-testid="stSidebar"] [data-testid="stExpander"]{
  border: 0px solid #F8FAFC;               /* gray-200 */
  border-radius: 4px;
  background: #F8FAFC;                      /* slate-50 */
  margin: 10px 0 12px 0;
}
[data-testid="stSidebar"] details > summary{
  padding: 10px 14px;                       /* header size like mock */
}
[data-testid="stSidebar"] [data-testid="stExpander"] div[role="group"]{
  padding: 8px 8px 12px 8px;                /* inner body padding */
}

/* 3) Uniform sidebar buttons */
[data-testid="stSidebar"] .stButton{ 
    margin: 2px 0; 
    border-color: #F8FAFC;
}

[data-testid="stSidebar"] .stButton > button{
  width: 100% !important;
  height: 44px;                             /* same height for all */
  border-radius: 4px;
  border: 0px solid #F8FAFC;                /* match card border */
  background: #FFFFFF;
  color: #111827;                            /* gray-900 text */
  font-weight: 500;
  justify-content: flex-start;               /* icon+text left aligned */
  gap: 8px;                                  /* space between emoji and text */
}
[data-testid="stSidebar"] .stButton > button:hover{
  background:#F9FAFB; border-color:#D1D5DB;
}

/* Divider: make it tight */
[data-testid="stSidebar"] hr { 
  margin: 2px 0 !important;          /* was large by default */
  border-color: #FFFFFF;
}

/* Caption under the title: reduce bottom space */
[data-testid="stSidebar"] .pluma-caption{
  margin-top: 0px;
  margin-bottom: 6px;                 /* add this */
}

/* Expander cards: tighter outer spacing + padding */
[data-testid="stSidebar"] [data-testid="stExpander"]{
  margin: 2px 0 4px 0 !important;     /* was 10px 0 12px 0 */
  border-radius: 2px;
}



[data-testid="stSidebar"] details > summary{
  padding: 2px 6px !important;       /* was 10px 14px */
}
[data-testid="stSidebar"] [data-testid="stExpander"] div[role="group"]{
  padding: 4px 6px 8px 6px !important;  /* was 8px 8px 12px 8px */
}

/* First expander right after the divider: remove extra top gap */
[data-testid="stSidebar"] .pluma-after-divider + [data-testid="stExpander"]{
  margin-top: 2px !important;
}
</style>
""", unsafe_allow_html=True)
pages.show_home()  # keep your existing header/banner etc.
st.session_state.setdefault("nav", "writer/style-writer")  # default route
st.session_state.setdefault("content", "")
st.session_state.setdefault("style", "")
st.session_state.setdefault("example", "")
st.session_state.setdefault("locals", {})  # you reference this

# ----------------------- Sidebar -------------------------
with st.sidebar:
    st.markdown('<div class="pluma-top">', unsafe_allow_html=True)
    c1, c2, c3 = st.columns([1, 3, 1])   # make center column wide
    with c2:
        st.image("img/pluma_logo.png", width=210)
    st.markdown('<div class="pluma-title">Pluma Writer</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="pluma-caption">BSP Style Writer • Edit, format, and validate content against your style guide.</div>',
        unsafe_allow_html=True
    )
    st.divider()
    st.markdown('<div class="pluma-after-divider"></div>', unsafe_allow_html=True)

    with st.expander("Chatbot", expanded=True):
        if st.button("💬  Open Chatbot", use_container_width=True):
            st.session_state["nav"] = "chatbot"

    with st.expander("Writer", expanded=True):
        if st.button("🖊️  Style Writer", use_container_width=True):
            st.session_state["nav"] = "writer/style-writer"
        if st.button("🔎  Style Reader", use_container_width=True):
            st.session_state["nav"] = "writer/style-reader"
        if st.button("📄  Generated Outputs", use_container_width=True):
            st.session_state["nav"] = "writer/outputs"

# ----------------------- Common header -------------------

route = st.session_state["nav"]


def show_style_reader():
    reader_page.render()
def show_outputs():
    outputs_page.render()


# ----------------------- Helpers -------------------------
def _register_pdf_font_if_available():
    try:
        font_path = os.path.join("assets", "DejaVuSans.ttf")
        if os.path.exists(font_path):
            pdfmetrics.registerFont(TTFont("DejaVuSans", font_path))
            return "DejaVuSans"
    except Exception:
        pass
    return "Helvetica"

def make_docx_bytes(text: str, title: str | None = None) -> bytes:
    doc = Document()
    if title:
        doc.add_heading(title, level=1)
    for block in text.replace("\r\n", "\n").split("\n\n"):
        p = doc.add_paragraph()
        for line in block.split("\n"):
            if line.strip():
                p.add_run(line)
            p.add_run("\n")
    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()

def make_pdf_bytes(text: str, title: str | None = None) -> bytes:
    font_name = _register_pdf_font_if_available()
    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm, title=title or "Rewrite",
        author="Style Writer",
    )
    styles = getSampleStyleSheet()
    base = styles["BodyText"]; base.fontName = font_name; base.fontSize = 11; base.leading = 14
    title_style = ParagraphStyle("Title", parent=styles["Heading1"], fontName=font_name, spaceAfter=12)
    story = []
    if title:
        story.append(Paragraph(title, title_style)); story.append(Spacer(1, 8))
    for block in text.replace("\r\n", "\n").split("\n\n"):
        block = block.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace("\n","<br/>")
        story.append(Paragraph(block, base)); story.append(Spacer(1, 6))
    doc.build(story)
    return buf.getvalue()

# ----------------------- Views ---------------------------
# def show_chainlit():
#     # Build Chainlit URL (append your own auth token if you have SSO)
#     params = {"user": st.session_state.get("user_id", "demo")}
#     CHAINLIT_BASE = "https://miniature-funicular-g4v7546vvr54cg6g-8000.app.github.dev"
#     url = f"{CHAINLIT_BASE}?{urlencode(params)}"
#     st.components.v1.iframe(url, height=600)

def show_chainlit():
    from urllib.parse import urlencode
    import streamlit as st

    # Build Chainlit URL
    params = {"user": st.session_state.get("user_id", "demo")}
    CHAINLIT_BASE = "https://miniature-funicular-g4v7546vvr54cg6g-8000.app.github.dev"
    url = f"{CHAINLIT_BASE}?{urlencode(params)}"

    # Add padding to the iframe container
    st.markdown(
        """
        <style>
        /* Add some top padding and remove extra margins */
        .chainlit-frame-container {
            padding-top: 18px;     /* 👈 Adjust this value (e.g. 12–24px) */
        }
        iframe[title="chainlit"] {
            border: none !important;
            border-radius: 8px;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

    # Display iframe inside padded container
    st.markdown('<div class="chainlit-frame-container">', unsafe_allow_html=True)
    st.components.v1.iframe(url, height=550, scrolling=True)
    st.markdown('</div>', unsafe_allow_html=True)


def show_style_writer():
    pages.show_sidebar()  # if this draws your Writer sidebar tools, keep it here
    st.header("📝Style Writer")

    # Content input
    st.session_state.content = st.text_area(
        ":blue[**Content Data:**]", st.session_state.content, 200
    )

    uploaded_files = st.file_uploader(
        ":blue[**Upload Content Files:**]",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True,
        help="Upload PDF, Word, or PowerPoint files",
    )

    # Extract text from uploads
    extracted_text = ""
    if uploaded_files:
        for uploaded_file in uploaded_files:
            file_type = uploaded_file.name.split(".")[-1].lower()
            if file_type == "pdf":
                pdf_reader = PdfReader(BytesIO(uploaded_file.read()))
                for page in pdf_reader.pages:
                    extracted_text += page.extract_text() + "\n"
            elif file_type == "docx":
                doc = Document(BytesIO(uploaded_file.read()))
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        extracted_text += paragraph.text + "\n"
            elif file_type == "pptx":
                prs = Presentation(BytesIO(uploaded_file.read()))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            extracted_text += shape.text + "\n"

    content_all = (
        st.session_state.content
        + "\n"
        + extracted_text.encode("ascii", errors="ignore").decode("ascii")
    )

    # Styles
    styles_data = utils.get_styles()
    style_options = [item["name"] for item in styles_data]
    selected_style = st.selectbox(":blue[**Select a Style:**]", options=style_options, index=None)

    if selected_style:
        filtered = next((item for item in styles_data if str(item["name"]) == selected_style), None)
        if filtered:
            st.session_state.style = filtered["style"]
            st.session_state.example = filtered["example"]
            st.session_state.styleId = selected_style

    # Guidelines
    guidelines = st.session_state.locals.get("relevant_guidelines", {})
    guidelines_summary = st.session_state.locals.get("guideline_summaries", {})
    selected_guidelines = []

    st.write(":blue[**Select Editorial Style Guides:**]")

    def render_guideline_checkbox(section_name: str, content: str, col_key_prefix: str):
        default_checked = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
        tooltip = guidelines_summary.get(section_name, None)
        if st.checkbox(
            section_name, value=default_checked,
            key=f"{col_key_prefix}_{section_name}", help=tooltip
        ):
            selected_guidelines.append(content)

    if guidelines:
        with st.container(border=True):
            col1, col2 = st.columns(2)
            items = list(guidelines.items()); mid = len(items) // 2
            with col1:
                for section_name, content in items[:mid]:
                    render_guideline_checkbox(section_name, content, "col1")
            with col2:
                for section_name, content in items[mid:]:
                    render_guideline_checkbox(section_name, content, "col2")
    else:
        st.warning("No guidelines available in the local data.")

    st.session_state.guidelines = "\n".join(selected_guidelines)

    # Note: Additional stylistic options are handled in the JS/React Writer UI.
    # The Streamlit prototype previously exposed quick checkboxes here; to avoid
    # duplicating behavior we keep that logic on the web frontend only.

    # Rewrite button
    disabled = (content_all == "" or st.session_state.style == "" or st.session_state.example == "")
    if st.button(":blue[**Rewrite Content**]", key="extract", disabled=disabled):
        with st.container(border=True):
            with st.spinner("Processing..."):
                st.markdown("### ✨ Rewritten Output")
                output = prompts.rewrite_content(content_all, False)
                utils.save_output(output, content_all)

                st.session_state["last_output"] = output
                ts = datetime.now().strftime("%Y%m%d-%H%M%S")
                style_id = (st.session_state.get("styleId") or "Style").replace(" ", "_")
                base_name = f"rewrite_{style_id}_{ts}"

                title_text = f"Rewrite • {st.session_state.get('styleId') or 'Selected Style'}"
                docx_bytes = make_docx_bytes(output, title=title_text)
                pdf_bytes = make_pdf_bytes(output, title=title_text)

                c1, c2 = st.columns(2)
                with c1:
                    st.download_button(
                        "⬇️ Download as DOCX", data=docx_bytes,
                        file_name=f"{base_name}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True,
                    )
                with c2:
                    st.download_button(
                        "⬇️ Download as PDF", data=pdf_bytes,
                        file_name=f"{base_name}.pdf", mime="application/pdf",
                        use_container_width=True,
                    )

# ----------------------- Router --------------------------
if route == "chatbot":
    show_chainlit()
elif route == "writer/style-writer":
    show_style_writer()
elif route == "writer/style-reader":
    show_style_reader()
elif route == "writer/outputs":
    show_outputs()
else:
    st.write("Not found.")
