# app/pages/reader.py
import streamlit as st
import app.pages as pages
import app.utils as utils
import app.prompts as prompts
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO

def render(key_prefix: str = "reader"):
    """Render the Style Reader page. key_prefix avoids widget key collisions."""
    # defaults
    st.session_state.setdefault("exampleText", "")
    st.session_state.setdefault("styleName", "")

    # header/side (keep if you want Reader to show same chrome)
    pages.show_home()
    pages.show_sidebar()

    st.header("🔍Style Reader")

    st.session_state.exampleText = st.text_area(
        ":blue[**Reference Style Example:**]",
        st.session_state.exampleText,
        height=200,
        key=f"{key_prefix}-example-text"
    )

    uploaded_files = st.file_uploader(
        ":blue[**Upload Example Files:**]",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True,
        help="Upload PDF, Word, or PowerPoint files",
        key=f"{key_prefix}-uploads"
    )

    st.session_state.styleName = st.text_input(
        ":blue[**Style Name:**]",
        st.session_state.styleName,
        max_chars=100,
        key=f"{key_prefix}-style-name"
    )

    # Extract text from uploaded files
    extracted_text = ""
    if uploaded_files:
        for uploaded_file in uploaded_files:
            ext = uploaded_file.name.split(".")[-1].lower()
            data = uploaded_file.read()

            if ext == "pdf":
                pdf_reader = PdfReader(BytesIO(data))
                for page in pdf_reader.pages:
                    extracted_text += (page.extract_text() or "") + "\n"

            elif ext == "docx":
                doc = Document(BytesIO(data))
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        extracted_text += paragraph.text + "\n"

            elif ext == "pptx":
                prs = Presentation(BytesIO(data))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        txt = getattr(shape, "text", "")
                        if txt and txt.strip():
                            extracted_text += txt + "\n"

    # Combine text area and extracted content
    combined_text = (
        (st.session_state.exampleText or "")
        + "\n"
        + (extracted_text.encode("ascii", errors="ignore").decode("ascii") if extracted_text else "")
    )

    if st.button(
        ":blue[**Extract Writing Style**]",
        key=f"{key_prefix}-extract",
        disabled=(combined_text.strip() == "" or st.session_state.styleName.strip() == "")
    ):
        with st.container(border=True):
            with st.spinner("Processing..."):
                if utils.check_style(st.session_state.styleName):
                    st.error(
                        f"Style name '{st.session_state.styleName}' already exists. "
                        "Please choose a different name."
                    )
                else:
                    style = prompts.extract_style(combined_text, False)
                    utils.save_style(style, combined_text)
                    st.success("✅ Style extracted and saved.")
