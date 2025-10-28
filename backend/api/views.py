from .repositories.factory import get_repo
from .services.prompts import extract_style as llm_extract_style, rewrite_content as llm_rewrite
from .services.config import LLM_LOCALS, load_locals
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework.permissions import AllowAny
from io import BytesIO
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# For chatbot
from django.shortcuts import get_object_or_404
from .models import ChatThread, ChatMessage, Attachment
from .services.chat_llm import run_chat

# For chat stream
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
from django.http import StreamingHttpResponse, HttpResponseBadRequest
from .services.llm_stream import stream_chat
from .services.sse import sse_format, sse_response
from .models import ChatThread, ChatMessage

# For Deep Research
import asyncio, threading, json
from queue import Queue, Empty
import time
import json, os
from .services.foundry_stream import stream_foundry_chat
from .services.foundry_stream import resolve_foundry
from .models import ChatFoundryThread

import logging
logger = logging.getLogger(__name__)

# Azure Foundry client libs
from azure.identity import DefaultAzureCredential
from azure.ai.agents import AgentsClient
from azure.ai.agents.models import (
    MessageAttachment,
    MessageInputTextBlock,
    MessageInputContentBlock,
    MessageImageFileParam,
    MessageInputImageFileBlock,
    FilePurpose,
    CodeInterpreterTool,
    AgentStreamEvent,
    MessageDeltaChunk,
    ThreadRun,
)
import tempfile
from pathlib import Path

repo = get_repo()

def _load_json_env(name: str) -> list[dict]:
    raw = os.getenv(name, "").strip()
    if not raw:
        return []
    try:
        return json.loads(raw)
    except Exception:
        return []


def _derive_title_from_text(text: str, max_len: int = 60) -> str:
    """Create a short, human-friendly title from a piece of text.
    Simple heuristic: take the first sentence or first N characters, collapse whitespace.
    """
    if not text:
        return ""
    s = text.strip()
    # pick up to the first sentence-ending punctuation
    import re
    m = re.search(r"([^.?!]+[.?!])", s)
    if m:
        s = m.group(1)
    # collapse whitespace
    s = " ".join(s.split())
    if len(s) > max_len:
        s = s[:max_len].rsplit(" ", 1)[0] + "..."
    return s


def extract_text_from_files(django_files):
    out = []
    for f in django_files:
        name = f.name.lower()
        data = f.read()
        if name.endswith(".pdf"):
            pdf = PdfReader(BytesIO(data))
            for p in pdf.pages:
                out.append((p.extract_text() or ""))
        elif name.endswith(".docx"):
            doc = Document(BytesIO(data))
            for p in doc.paragraphs:
                if p.text.strip(): out.append(p.text)
        elif name.endswith(".pptx"):
            prs = Presentation(BytesIO(data))
            for s in prs.slides:
                for sh in s.shapes:
                    txt = getattr(sh, "text", "")
                    if txt and txt.strip(): out.append(txt)
    return ("\n".join(out)).encode("ascii","ignore").decode("ascii")


class StylesAPI(APIView):
    def get(self, _):
        # Try to honor an Azure-authenticated user id header like the old Streamlit app.
        # Frontend (browser) may not send this header; when present we pass it to the repo
        # so CosmosRepo can filter by user_id.
        try:
            user_id = _.headers.get("X-MS-CLIENT-PRINCIPAL-ID") or _.META.get("HTTP_X_MS_CLIENT_PRINCIPAL_ID")
        except Exception:
            user_id = None
        return Response(repo.list_styles(user_id=user_id))

    def post(self, request):
        name = (request.data.get("name") or "").strip()
        if not name:
            return Response({"detail":"name required"}, status=400)
        data = repo.create_or_update_style(name=name,
                                           style=request.data.get("style",""),
                                           example=request.data.get("example",""))
        return Response({"id": data["id"], "name": data["name"]})

class StyleDetailAPI(APIView):
    def delete(self, request, style_id):
        repo.delete_style(str(style_id))
        return Response(status=204)

class OutputsAPI(APIView):
    def get(self, _):
        return Response(repo.list_outputs())


class LocalsAPI(APIView):
    """Return the local_data.json contents used for prompts (LLM_LOCALS)
    GET /api/locals/
    """
    def get(self, request):
        # Only return specific safe keys for the frontend. Read fresh to pick up file changes.
        try:
            locals_data = load_locals() or {}
        except Exception:
            locals_data = LLM_LOCALS or {}
        allowed = ["relevant_guidelines", "guideline_summaries"]
        out = {k: locals_data.get(k, {}) for k in allowed}
        return Response(out)


class OutputDownloadAPI(APIView):
    """Download a saved output as PDF or DOCX.
    GET /api/outputs/<output_id>/download/?format=pdf|docx
    """
    def get(self, request, output_id: str):
        fmt = (request.query_params.get("format") or "pdf").lower()
        if fmt not in ("pdf", "docx"):
            return Response({"detail": "format must be 'pdf' or 'docx'"}, status=400)

        data = repo.get_output(str(output_id))
        if not data:
            return Response({"detail": "output not found"}, status=404)

        text = (data.get("output") or "")
        from io import BytesIO
        if fmt == "docx":
            from docx import Document
            buf = BytesIO()
            doc = Document()
            # preserve paragraphs
            for line in text.splitlines():
                doc.add_paragraph(line)
            doc.save(buf)
            buf.seek(0)
            resp = Response(buf.read(), status=200, content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            resp["Content-Disposition"] = f"attachment; filename=output_{output_id}.docx"
            return resp

        # PDF
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        buf = BytesIO()
        c = canvas.Canvas(buf, pagesize=letter)
        width, height = letter
        margin = 50
        y = height - margin
        line_height = 12
        for paragraph in text.split("\n\n"):
            for line in paragraph.splitlines():
                # simple wrap at ~90 chars
                import textwrap
                wrapped = textwrap.wrap(line, width=90) or [""]
                for wl in wrapped:
                    if y < margin + line_height:
                        c.showPage()
                        y = height - margin
                    c.setFont("Helvetica", 10)
                    c.drawString(margin, y, wl)
                    y -= line_height
            y -= line_height  # paragraph gap
        c.save()
        buf.seek(0)
        resp = Response(buf.read(), status=200, content_type="application/pdf")
        resp["Content-Disposition"] = f"attachment; filename=output_{output_id}.pdf"
        return resp

class RewriteAPI(APIView):
    def post(self, request):
        content = request.data.get("content","")
        if not content.strip():
            return Response({"detail":"content empty"}, status=400)
        style      = request.data.get("style","")
        example    = request.data.get("example","")
        guidelines = request.data.get("guidelines","")
        style_id   = request.data.get("styleId","Style")

        from .services.prompts import rewrite_content as llm_rewrite
        rewritten = llm_rewrite(content_all=content, style=style, guidelines=guidelines, example=example)
        saved = repo.save_output(style_name=style_id, input_text=content, output_text=rewritten)
        return Response({"output": rewritten, "output_id": saved["id"]})

from .services.prompts import extract_style as llm_extract_style

class ExtractStyleAPI(APIView):
    def post(self, request):
        example_text = request.data.get("exampleText", "")
        files = request.FILES.getlist("files")
        combined = (example_text + "\n" + extract_text_from_files(files)).strip()
        if not combined:
            return Response({"detail": "empty input"}, status=status.HTTP_400_BAD_REQUEST)

        # Some Azure deployments only accept the default model temperature (1.0).
        # Use 1.0 here to avoid invalid_request_error from the provider.
        style_text = llm_extract_style(combined_text=combined, temperature=1.0)
        return Response({"style": style_text})


class ChatStartAPI(APIView):
    def post(self, request):
        user_id = request.headers.get("X-MS-CLIENT-PRINCIPAL-ID", "")  # Azure header if present
        title = (request.data.get("title") or "").strip()
        th = ChatThread.objects.create(title=title, user_id=user_id)
        # seed with a system message if you want to customize per-thread
        ChatMessage.objects.create(thread=th, role="system", content="")
        # create an initial assistant greeting so new chats show a welcoming message
        try:
            ChatMessage.objects.create(thread=th, role="assistant", content="Hello, how can I assist you today?")
        except Exception:
            # non-fatal: don't fail thread creation if greeting save fails
            pass
        return Response({"thread_id": th.id}, status=201)

# backend/api/views.py
class ChatHistoryAPI(APIView):
    def get(self, request):
        # Support either `thread_id` or `thread` query param. If missing/invalid, return empty history
        tid = request.query_params.get("thread_id") or request.query_params.get("thread")
        try:
            if not tid:
                return Response({"thread_id": None, "messages": []})
            thread_id = int(tid)
        except Exception:
            return Response({"thread_id": None, "messages": []})

        try:
            th = get_object_or_404(ChatThread, id=thread_id)
        except Exception:
            return Response({"thread_id": None, "messages": []})

        msgs = [
            {
                "role": m.role,
                "content": m.content,
                "created_at": m.created_at.isoformat(),
                "attachments": [
                    {
                        "id": a.id,
                        "filename": a.filename,
                        "content_type": a.content_type,
                        "blob_url": a.blob_url,
                        "foundry_file_id": a.foundry_file_id,
                        "created_at": a.created_at.isoformat(),
                    }
                    for a in m.attachments.order_by("created_at").all()
                ],
            }
            for m in th.messages.order_by("created_at")
        ]
        return Response({"thread_id": th.id, "messages": msgs})


class ChatMessageAPI(APIView):
    def post(self, request):
        """
        body: { thread_id, content, temperature? }
        """
        data = request.data
        thread_id = data.get("thread_id")
        content = (data.get("content") or "").strip()
        temperature = float(data.get("temperature") or 0.7)

        if not thread_id or not content:
            return Response({"detail": "thread_id and content required"}, status=400)

        th = get_object_or_404(ChatThread, id=thread_id)

        # Save user message
        ChatMessage.objects.create(thread=th, role="user", content=content)

        # If the thread has no meaningful title yet, derive one from the user's first message
        try:
            cur_title = (th.title or "").strip()
            if not cur_title or cur_title.lower().startswith("new") or cur_title.lower().startswith("chat"):
                new_title = _derive_title_from_text(content) or cur_title
                if new_title and new_title != cur_title:
                    th.title = new_title
                    th.save()
        except Exception:
            # non-fatal; don't block the request
            pass

        # Build history for LLM
        history = [
            {"role": m.role, "content": m.content}
            for m in th.messages.order_by("created_at")
            if m.content.strip()
        ]

        # Call Azure LLM
        assistant_text = run_chat(history, temperature=temperature)

        # Save assistant message
        ChatMessage.objects.create(thread=th, role="assistant", content=assistant_text or "")

        return Response({"output": assistant_text})


class ChatUploadAPI(APIView):
    """
    POST /api/chat/upload/
    Accepts multipart/form-data files[] and returns extracted text.
    Body: files[]
    Response: { text: str }
    """
    def post(self, request):
        files = request.FILES.getlist("files")
        if not files:
            return Response({"text": ""})
        try:
            text = extract_text_from_files(files)
            return Response({"text": text})
        except Exception:
            return Response({"text": ""}, status=500)


class ChatUploadFoundryAPI(APIView):
    """
    POST /api/chat/upload-foundry/
    Accepts multipart/form-data files[], thread_id, optional content, model_deployment and mode.
    Uploads files to Azure Foundry, creates a Foundry message with attachments and streams the run as SSE.
    """
    permission_classes = [AllowAny]

    def post(self, request):
        # For multipart requests, use POST dict
        data = request.POST if request.content_type.startswith("multipart/") else request.data
        thread_id = data.get("thread_id") or data.get("thread")
        user_text = (data.get("content") or "").strip()
        deployment = data.get("model_deployment") or data.get("deployment") or data.get("model")
        mode = (data.get("mode") or "work").lower()

        files = request.FILES.getlist("files")

        if not thread_id:
            return Response({"detail": "thread_id required"}, status=400)
        if not deployment:
            return Response({"detail": "model_deployment required"}, status=400)

        # resolve foundry endpoint and agent id
        try:
            resolved = resolve_foundry(model_deployment=deployment, mode=mode)
            endpoint = resolved["endpoint"]
            agent_id = resolved["model_id"]
        except Exception as e:
            return Response({"detail": f"foundry resolve failed: {e}"}, status=500)

        # prepare credential and client
        cred = DefaultAzureCredential(exclude_shared_token_cache_credential=True)
        client = AgentsClient(endpoint=endpoint, credential=cred)

        # ensure mapping from ChatThread -> foundry thread
        try:
            th = get_object_or_404(ChatThread, id=int(thread_id))
        except Exception:
            return Response({"detail": "thread not found"}, status=404)

        f_thread_id = None
        try:
            mapping = ChatFoundryThread.objects.filter(thread=th).first()
            if mapping and mapping.foundry_thread_id:
                f_thread_id = mapping.foundry_thread_id
        except Exception:
            mapping = None

        if not f_thread_id:
            t = client.threads.create()
            f_thread_id = t.id
            try:
                ChatFoundryThread.objects.update_or_create(thread=th, defaults={"foundry_thread_id": f_thread_id})
            except Exception:
                pass

        # helper to emit SSE frames
        def sse_event(name: str, data_str: str):
            yield f"event: {name}\n"
            for line in data_str.splitlines() or [""]:
                yield f"data: {line}\n"
            yield "\n"

        def gen():
            yield from sse_event("ready", "{}")

            # persist user message early so history reflects it (include attachment markers)
            user_msg = None
            try:
                if user_text or files:
                    marker_text = ""
                    try:
                        names = [f.name for f in files]
                        if names:
                            marker_text = "\n\n" + "\n".join([f"<file_name:{n}>" for n in names])
                    except Exception:
                        marker_text = ""
                    user_msg = ChatMessage.objects.create(thread=th, role="user", content=(user_text or "") + marker_text)
            except Exception:
                user_msg = None

            # upload files and prepare attachments/content blocks
            content_blocks = [MessageInputTextBlock(text=user_text)]
            attachments = []
            tmp_paths = []
            try:
                for f in files:
                    # write to a temp file
                    suffix = Path(f.name).suffix or ""
                    tf = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
                    tf.write(f.read())
                    tf.flush()
                    tf.close()
                    tmp_paths.append(tf.name)

                    uploaded = client.files.upload_and_poll(file_path=tf.name, purpose=FilePurpose.AGENTS)
                    attachments.append(MessageAttachment(file_id=uploaded.id, tools=CodeInterpreterTool().definitions))

                    # Also upload to blob storage (if configured) and persist an Attachment record
                    try:
                        from azure.storage.blob import BlobServiceClient
                        conn = os.getenv("APP_AZURE_STORAGE_CONNECTION_STRING") or os.getenv("APP_AZURE_STORAGE_CONNECTION")
                        blob_url = ""
                        if conn:
                            bsc = BlobServiceClient.from_connection_string(conn)
                            container = os.getenv("BUCKET_NAME") or os.getenv("AZURE_STORAGE_CONTAINER") or "chats"
                            blob_name = f"{int(time.time()*1000)}_{Path(tf.name).name}"
                            blob_client = bsc.get_blob_client(container=container, blob=blob_name)
                            with open(tf.name, "rb") as fh:
                                blob_client.upload_blob(fh, overwrite=True)
                            # Try to compute a public URL using DEV_AZURE_BLOB_ENDPOINT if provided
                            dev_endpoint = os.getenv("DEV_AZURE_BLOB_ENDPOINT")
                            if dev_endpoint:
                                blob_url = f"{dev_endpoint.rstrip('/')}" + f"/{container}/{blob_name}"
                            else:
                                try:
                                    blob_url = blob_client.url
                                except Exception:
                                    blob_url = ""
                        else:
                            blob_url = ""

                        # Create Attachment record linked to the persisted user message (if present)
                        if user_msg:
                            Attachment.objects.create(
                                message=user_msg,
                                filename=f.name,
                                content_type=(f.content_type or ""),
                                blob_url=blob_url or "",
                                foundry_file_id=(getattr(uploaded, "id", "") or ""),
                            )
                    except Exception:
                        # don't fail the stream if blob upload/attachment save fails
                        pass

                    # if image, add image content block
                    if (f.content_type or "").startswith("image/"):
                        file_param = MessageImageFileParam(file_id=uploaded.id, detail="high")
                        content_blocks.append(MessageInputImageFileBlock(image_file=file_param))
            except Exception as e:
                # cleanup tmp files
                for p in tmp_paths:
                    try:
                        Path(p).unlink()
                    except Exception:
                        pass
                yield from sse_event("error", json.dumps({"detail": str(e)}))
                yield from sse_event("done", "{}")
                return

            # create the message in foundry thread
            try:
                client.messages.create(thread_id=f_thread_id, role="user", content=content_blocks, attachments=attachments)
            except Exception as e:
                yield from sse_event("error", json.dumps({"detail": str(e)}))
                yield from sse_event("done", "{}")
                return

            # stream run
            tokens: list[str] = []
            try:
                with client.runs.stream(thread_id=f_thread_id, agent_id=agent_id) as stream:
                    for event_type, event_data, _ in stream:
                        if isinstance(event_data, MessageDeltaChunk):
                            token = event_data.text or ""
                            if token:
                                tokens.append(token)
                                yield from sse_event("token", token)
                        elif isinstance(event_data, ThreadRun):
                            if event_data.status == "failed":
                                raise RuntimeError(str(event_data.last_error))
                        elif event_type == AgentStreamEvent.ERROR:
                            raise RuntimeError(str(event_data))

                # merge tokens into assistant text
                def _merge_tokens(ts: list[str]) -> str:
                    if not ts:
                        return ""
                    out = ts[0]
                    for t in ts[1:]:
                        if out and out[-1].isalnum() and t and t[0].isalnum():
                            out += " " + t
                        else:
                            out += t
                    return out

                assistant_text = _merge_tokens(tokens)
                try:
                    ChatMessage.objects.create(thread=th, role="assistant", content=assistant_text)
                except Exception:
                    pass

                # update thread title if needed
                try:
                    cur_title = (th.title or "").strip()
                    if not cur_title or cur_title.lower().startswith("new") or cur_title.lower().startswith("chat"):
                        new_title = _derive_title_from_text(assistant_text) or cur_title
                        if new_title and new_title != cur_title:
                            th.title = new_title
                            th.save()
                except Exception:
                    pass

                yield from sse_event("done", json.dumps({"ok": True}))
            except Exception as e:
                yield from sse_event("error", json.dumps({"detail": str(e)}))
                yield from sse_event("done", "{}")
            finally:
                # cleanup tmp files
                for p in tmp_paths:
                    try:
                        Path(p).unlink()
                    except Exception:
                        pass

        resp = StreamingHttpResponse(gen(), content_type="text/event-stream; charset=utf-8")
        resp["Cache-Control"] = "no-cache"
        resp["X-Accel-Buffering"] = "no"
        return resp

class ChatModelsAPI(APIView):
    """
    GET /api/chat/models/  ->  { models: LLM_CONFIG, workweb: LLM_WORKWEB }
    Used by frontend to populate the model selector and Work/Web toggle.
    """
    def get(self, _):
        return Response({
            "models": _load_json_env("LLM_CONFIG"),
            "workweb": _load_json_env("LLM_WORKWEB"),
        })


class ChatThreadsAPI(APIView):
    """List chat threads with metadata for the sidebar.
    GET /api/chat/threads/
    Returns: [{thread_id, title, last_message, last_updated, created_at}]
    """
    def get(self, request):
        threads = ChatThread.objects.order_by("-created_at").all()[:50]
        out = []
        for th in threads:
            last = th.messages.order_by("-created_at").first()
            out.append({
                "thread_id": th.id,
                "title": th.title or (last.content[:60] if last else "New chat"),
                "last_message": last.content if last else "",
                "last_updated": last.created_at.isoformat() if last else th.created_at.isoformat(),
                "created_at": th.created_at.isoformat(),
            })
        return Response(out)


class ChatRenameAPI(APIView):
    """Rename a chat thread title.
    POST /api/chat/rename/ body: { thread_id, title }
    """
    def post(self, request):
        thread_id = request.data.get("thread_id")
        title = (request.data.get("title") or "").strip()
        if not thread_id:
            return Response({"detail": "thread_id required"}, status=400)
        th = get_object_or_404(ChatThread, id=thread_id)
        th.title = title
        th.save()
        return Response({"thread_id": th.id, "title": th.title})


@method_decorator(csrf_exempt, name="dispatch")
class ChatStreamAPI(APIView):
    """
    Minimal Server-Sent Events endpoint.
    Streams back tokens so your React UI can render progressively.
    Replace the generator body with your Foundry streaming later.
    """
    permission_classes = [AllowAny]

    def post(self, request, *args, **kwargs):
        # Payload your UI sends:
        # { thread_id, content, provider: "foundry", model_deployment, mode }
        try:
            data = request.data if request.content_type != "application/x-www-form-urlencoded" else request.POST
            user_text = (data.get("content") or "").strip()
            thread_id = data.get("thread_id") or data.get("thread")
            provider = (data.get("provider") or "").lower()
            deployment = data.get("deployment") or data.get("model_deployment") or data.get("model")
            mode = (data.get("mode") or "work").lower()
        except Exception:
            user_text = ""

        def sse_event(name: str, data_str: str):
            yield f"event: {name}\n"
            # each line of data must be prefixed with "data: "
            for line in data_str.splitlines() or [""]:
                yield f"data: {line}\n"
            yield "\n"

        # try to resolve thread object early so we can persist the user's message
        th_obj = None
        if thread_id:
            try:
                th_obj = get_object_or_404(ChatThread, id=int(thread_id))
            except Exception:
                th_obj = None

        def gen():
            # tell client we’re ready
            yield from sse_event("ready", "{}")

            # If client requested Foundry provider, use the streaming implementation
            if provider == "foundry":
                if not thread_id:
                    yield from sse_event("error", json.dumps({"detail": "thread_id required for streaming"}))
                    yield from sse_event("done", "{}")
                    return

                try:
                    # stream_foundry_chat yields token strings (not SSE frames)
                    # ensure we have the thread and persist the user's message
                    th = th_obj or get_object_or_404(ChatThread, id=int(thread_id))
                    try:
                        if user_text:
                            ChatMessage.objects.create(thread=th, role="user", content=user_text)
                    except Exception:
                        # non-fatal: continue streaming even if save fails
                        pass
                    tokens = []
                    for token in stream_foundry_chat(thread_db_id=int(thread_id), user_text=user_text, model_deployment=deployment or "", mode=mode):
                        # each token is a string chunk
                        tokens.append(token)
                        yield from sse_event("token", token)
                    # persist final assistant message (merge tokens with spacing rules)
                    try:
                        def _merge_tokens(ts: list[str]) -> str:
                            if not ts:
                                return ""
                            out = ts[0]
                            for t in ts[1:]:
                                if out and out[-1].isalnum() and t and t[0].isalnum():
                                    out += " " + t
                                else:
                                    out += t
                            return out

                        assistant_text = _merge_tokens(tokens)
                        ChatMessage.objects.create(thread=th, role="assistant", content=assistant_text)
                        # update thread title if it's still default/empty using assistant text as fallback
                        try:
                            cur_title = (th.title or "").strip()
                            if not cur_title or cur_title.lower().startswith("new") or cur_title.lower().startswith("chat"):
                                new_title = _derive_title_from_text(assistant_text) or cur_title
                                if new_title and new_title != cur_title:
                                    th.title = new_title
                                    th.save()
                        except Exception:
                            pass
                    except Exception:
                        # don't crash the stream if DB save fails
                        pass
                    yield from sse_event("done", json.dumps({"ok": True}))
                    return
                except Exception as e:
                    try:
                        detail = {"detail": str(e)}
                        yield from sse_event("error", json.dumps(detail))
                    except Exception:
                        yield from sse_event("error", json.dumps({"detail": "streaming failed"}))
                    yield from sse_event("done", "{}")
                    return

            # Fallback demo stream when no provider or not foundry
            if not user_text:
                yield from sse_event("token", "Hello! 👋")
                yield from sse_event("done", "{}")
                return

            # For non-foundry fallback streaming, persist the user's message if we have a thread
            try:
                if th_obj and user_text:
                    try:
                        ChatMessage.objects.create(thread=th_obj, role="user", content=user_text)
                    except Exception:
                        pass
            except Exception:
                pass

            for word in (user_text.split() or []):
                yield from sse_event("token", word + " ")
                time.sleep(0.03)  # small drip so you see streaming

            # Example: send an extra completion suffix
            yield from sse_event("token", "\n\n(placeholder stream — wire Foundry here)")
            yield from sse_event("done", json.dumps({"ok": True}))

        resp = StreamingHttpResponse(gen(), content_type="text/event-stream; charset=utf-8")
        # important for proxies/buffers
        resp["Cache-Control"] = "no-cache"
        resp["X-Accel-Buffering"] = "no"
        return resp


@method_decorator(csrf_exempt, name="dispatch")
class ResearchStreamAPI(APIView):
    """
    POST /api/research/stream
    body: { topic: str }
    Streams SSE events from deep research pipeline.
    """
    def post(self, request):
        topic = (request.data.get("topic") or "").strip()
        if not topic:
            return Response({"detail": "topic required"}, status=400)

        q: Queue[bytes] = Queue(maxsize=100)

        # bridge: the notify callback enqueues SSE frames
        async def notify(event: str, payload: dict):
            try:
                frame = sse_format(json.dumps(payload), event=event)
                q.put(frame, timeout=2)
            except Exception:
                pass

        # run the async pipeline in a dedicated thread with its own loop
        def runner():
            # create and set a fresh event loop for this thread BEFORE importing the pipeline
            loop = asyncio.new_event_loop()
            # bind this loop to the current thread
            asyncio.set_event_loop(loop)

            try:
                # now import the pipeline so any async clients it constructs bind to this loop
                try:
                    from .deep_research import pipeline as dr
                except Exception as imp_e:
                    logger.exception("deep_research import failed")
                    try:
                        q.put(sse_format(json.dumps({"detail": f"deep_research import failed: {imp_e}"}), event="error"), timeout=2)
                    except Exception:
                        pass
                    q.put(b"__CLOSE__")
                    return

                # Emit a quick debug frame to confirm the import happened and the
                # runner thread is active. This helps determine if the pipeline is
                # being executed even when the LLM doesn't emit 'thinking' notes.
                try:
                    q.put(sse_format(json.dumps({"detail": "deep_research imported"}), event="debug"))
                except Exception:
                    pass

                async def main_and_cleanup():
                    try:
                        result = await dr.run_deep_research(topic, notify=notify)
                        # final "done" event with summary
                        q.put(sse_format(json.dumps({"summary": result}), event="done"))
                    except Exception as e:
                        logger.exception("deep_research runtime error")
                        try:
                            q.put(sse_format(json.dumps({"detail": str(e)}), event="error"), timeout=2)
                        except Exception:
                            pass
                    finally:
                        # attempt to gracefully close/cleanup async clients created by the pipeline module
                        try:
                            # common pattern: models/clients exposing an aclose coroutine
                            ds = getattr(dr, "deep_seek_model", None)
                            if ds is not None:
                                aclose = getattr(ds, "aclose", None) or getattr(ds, "close", None)
                                if aclose is not None:
                                    if asyncio.iscoroutinefunction(aclose):
                                        await aclose()
                                    else:
                                        try:
                                            aclose()
                                        except Exception:
                                            pass
                        except Exception:
                            logger.exception("deep_research cleanup failed")
                        # sentinel to close the stream
                        q.put(b"__CLOSE__")

                # Run the pipeline and cleanup on this loop
                try:
                    # notify via queue that we're about to start the async run
                    try:
                        q.put(sse_format(json.dumps({"detail": "starting deep_research run"}), event="debug"))
                    except Exception:
                        pass
                    loop.run_until_complete(main_and_cleanup())
                except Exception as run_e:
                    logger.exception("deep_research async run failed")
                    try:
                        q.put(sse_format(json.dumps({"detail": f"async run failed: {run_e}"}), event="error"), timeout=2)
                    except Exception:
                        pass

                # After the main run completes, try to cancel any lingering tasks and
                # allow async generators to finish before closing the loop.
                try:
                    pending = asyncio.all_tasks(loop)
                    for t in pending:
                        try:
                            t.cancel()
                        except Exception:
                            pass
                    if pending:
                        try:
                            loop.run_until_complete(asyncio.gather(*pending, return_exceptions=True))
                        except Exception:
                            pass
                    try:
                        loop.run_until_complete(loop.shutdown_asyncgens())
                    except Exception:
                        pass
                except Exception:
                    logger.exception("error during deep_research loop shutdown")
            finally:
                try:
                    loop.close()
                except Exception:
                    pass

        threading.Thread(target=runner, daemon=True).start()

        def gen():
            # immediate ack
            yield sse_format("started", event="ready")
            while True:
                try:
                    item = q.get(timeout=30)
                    if item == b"__CLOSE__":
                        break
                    yield item
                except Empty:
                    # keep-alive
                    yield sse_format("ping", event="keepalive")

        return sse_response(gen())


@method_decorator(csrf_exempt, name="dispatch")
class ReasoningStreamAPI(APIView):
    """
    POST /api/reasoning/stream
    body: { query: str, provider?: str, model_deployment?: str, mode?: 'work'|'web' }
    Streams minimal SSE status events and returns final markdown.
    """
    def post(self, request):
        query = (request.data.get("query") or "").strip()
        provider = (request.data.get("provider") or "").strip().lower() or None
        model_deployment = (request.data.get("model_deployment") or request.data.get("deployment") or request.data.get("model") or "").strip() or None
        mode = (request.data.get("mode") or "").strip().lower() or None
        if not query:
            return Response({"detail": "query required"}, status=400)

        q: Queue[bytes] = Queue(maxsize=100)

        async def notify(payload: dict):
            try:
                frame = sse_format(json.dumps(payload), event=str(payload.get("event") or "message"))
                q.put(frame, timeout=2)
            except Exception:
                pass

        def runner():
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)

            try:
                try:
                    from .reasoning import reasoning as rsn
                except Exception as imp_e:
                    logger.exception("reasoning import failed")
                    try:
                        q.put(sse_format(json.dumps({"detail": f"reasoning import failed: {imp_e}"}), event="error"), timeout=2)
                    except Exception:
                        pass
                    q.put(b"__CLOSE__")
                    return

                async def main_and_cleanup():
                    try:
                        markdown = await rsn.run_reasoning(query, notify=notify, provider=provider, model_deployment=model_deployment, mode=mode)
                        q.put(sse_format(json.dumps({"markdown": markdown}), event="done"))
                    except Exception as e:
                        logger.exception("reasoning runtime error")
                        try:
                            q.put(sse_format(json.dumps({"detail": str(e)}), event="error"), timeout=2)
                        except Exception:
                            pass
                    finally:
                        q.put(b"__CLOSE__")

                try:
                    loop.run_until_complete(main_and_cleanup())
                except Exception as run_e:
                    logger.exception("reasoning async run failed")
                    try:
                        q.put(sse_format(json.dumps({"detail": f"async run failed: {run_e}"}), event="error"), timeout=2)
                    except Exception:
                        pass
                try:
                    pending = asyncio.all_tasks(loop)
                    for t in pending:
                        try:
                            t.cancel()
                        except Exception:
                            pass
                    if pending:
                        try:
                            loop.run_until_complete(asyncio.gather(*pending, return_exceptions=True))
                        except Exception:
                            pass
                    try:
                        loop.run_until_complete(loop.shutdown_asyncgens())
                    except Exception:
                        pass
                except Exception:
                    logger.exception("error during reasoning loop shutdown")
            finally:
                try:
                    loop.close()
                except Exception:
                    pass

        threading.Thread(target=runner, daemon=True).start()

        def gen():
            yield sse_format("started", event="ready")
            while True:
                try:
                    item = q.get(timeout=30)
                    if item == b"__CLOSE__":
                        break
                    yield item
                except Empty:
                    yield sse_format("ping", event="keepalive")

        return sse_response(gen())